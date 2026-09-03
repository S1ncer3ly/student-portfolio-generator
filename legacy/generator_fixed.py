import streamlit as st
import pandas as pd
from pptx import Presentation
import os
from PIL import Image
import pillow_heif
import shutil
from pathlib import Path

# ================= PAGE CONFIG =================
st.set_page_config(page_title="Student Portfolio Generator v2.0", page_icon="🎓", layout="wide")

# BIG VERSION LABEL TO VERIFY THE VERSION
st.title("🎓 Student Portfolio Generator v2.0")
st.success("🚀 FIXED COORDINATE MODE ACTIVE - Photos will now stay in place!")
st.markdown("Automate your student presentations. **Strict Naming Mode Enabled.**")

# ================= HELPER FUNCTIONS =================
def convert_heic_to_jpg(heic_path, output_path):
    """Converts HEIC photos to JPEG for PPTX compatibility."""
    try:
        heif_file = pillow_heif.read_heif(heic_path)
        image = Image.frombytes(
            heif_file.mode, heif_file.size, heif_file.data, "raw", heif_file.mode, heif_file.stride
        )
        image.save(output_path, "JPEG")
        return True
    except Exception as e:
        st.error(f"Error converting {heic_path}: {e}")
        return False

def replace_text_recursive(shapes, replacements):
    """Deep search and replace for tags like {{NAME}} inside groups."""
    for shape in shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in replacements.items():
                        if key in run.text:
                            run.text = run.text.replace(key, str(value))
        if shape.shape_type == 6: # Group Shape
            replace_text_recursive(shape.shapes, replacements)

# ================= SIDEBAR (GLOBAL SETTINGS) =================
with st.sidebar:
    st.header("⚙️ Global Settings")
    root_path = st.text_input("Root Photo Folder Path", placeholder="/Users/name/Desktop/Grade 1")
    global_theme = st.text_input("Global Theme Name", value="Heroes, Role Models & Changemakers")
    uploaded_template = st.file_uploader("Upload PPTX Template", type="pptx")
    st.divider()
    st.warning("⚠️ **Strict Naming Required:** Photos must be named exactly `STUDENT NAME_W1.heic` etc.")

# ================= MAIN WORKSPACE =================
tab1, tab2, tab3 = st.tabs(["📥 Data Import", "✅ Pre-Flight Check", "🚀 Generation"])

with tab1:
    st.subheader("Step 1: Import Student List")
    uploaded_data = st.file_uploader("Upload Student CSV or Excel", type=["csv", "xlsx"])
    df = None
    if uploaded_data:
        try:
            if uploaded_data.name.endswith('.csv'):
                df = pd.read_csv(uploaded_data)
            else:
                df = pd.read_excel(uploaded_data)
            if df.empty:
                st.error("The uploaded file is empty.")
                df = None
            else:
                st.write("### Preview of Student Data")
                st.dataframe(df, use_container_width=True)
                st.divider()
                st.write("### Column Mapping")
                col1, col2 = st.columns(2)
                with col1:
                    name_col = st.selectbox("Which column is the Student Name?", options=df.columns)
                    class_col = st.selectbox("Which column is the Class?", options=df.columns)
                with col2:
                    section_col = st.selectbox("Which column is the Section?", options=df.columns)
                    theme_col = st.selectbox("Which column is the Theme? (Optional)", options=[None] + list(df.columns))
        except Exception as e:
            st.error(f"An error occurred: {e}")

with tab2:
    st.subheader("Step 2: Verify Photos")
    if st.button("🔍 Scan Folders for Photos"):
        if not root_path or uploaded_data is None or df is None:
            st.warning("Please provide the Root Folder Path and upload a valid Student List first!")
        else:
            results = []
            for _, row in df.iterrows():
                name = str(row[name_col])
                status = {"Student": name, "Week 1": "❌", "Week 2": "❌", "Week 3": "❌", "Week 4": "❌", "Overall": "🔴"}
                all_found = True
                for w in range(1, 5):
                    expected_filename = f"{name}_W{w}.heic"
                    folder = os.path.join(root_path, f"Week {w}")
                    if os.path.exists(folder) and expected_filename in os.listdir(folder):
                        status[f"Week {w}"] = "✅"
                    else:
                        all_found = False
                status["Overall"] = "🟢" if all_found else "🟡"
                results.append(status)
            st.table(pd.DataFrame(results))

with tab3:
    st.subheader("Step 3: Generate Presentations")
    if st.button("🔥 GENERATE ALL PRESENTATIONS", type="primary", use_container_width=True):
        if not uploaded_template or not uploaded_data or not root_path or df is None:
            st.error("Missing requirements!")
        else:
            with open("temp_template.pptx", "wb") as f:
                f.write(uploaded_template.getbuffer())

            output_folder = os.path.join(root_path, "Finished_PPTs")
            if not os.path.exists(output_folder):
                os.makedirs(output_folder)

            progress_bar = st.progress(0)
            status_text = st.empty()

            # FIXED COORDINATES extracted from template (L, T, W, H)
            # These are in EMUs (English Metric Units)
            COORD_MAP = {
                3: (1603584, 3260950, 9994200, 5740500), # Slide 3 / Week 1
                4: (1603584, 3260950, 9994200, 5740500), # Slide 4 / Week 2
                5: (1603584, 3260950, 9994200, 5740500), # Slide 5 / Week 3
                6: (1603584, 3260950, 9994200, 5740500), # Slide 6 / Week 4
            }

            for index, row in df.iterrows():
                name = str(row[name_col])
                s_class = str(row[class_col])
                section = str(row[section_col])
                theme = str(row[theme_col]) if theme_col and pd.notna(row[theme_col]) else global_theme
                status_text.text(f"Processing {index+1}/{len(df)}: {name}...")
                try:
                    prs = Presentation("temp_template.pptx")
                    replacements = {"{{NAME}}": name, "{{CLASS}}": s_class, "{{SECTION}}": section, "{{THEME}}": theme}
                    replace_text_recursive(prs.slides[1].shapes, replacements)

                    for w in range(1, 5):
                        slide_idx = w + 2
                        slide = prs.slides[slide_idx - 1]

                        # Remove ANY existing shapes at that location to avoid overlaps
                        for shape in list(slide.shapes):
                            if shape.is_placeholder and shape.placeholder_format.type == 18:
                                element = shape._element
                                element.getparent().remove(element)

                        expected_filename = f"{name}_W{w}.heic"
                        photo_path = os.path.join(root_path, f"Week {w}", expected_filename)
                        if os.path.exists(photo_path):
                            temp_jpg = f"temp_{index}_{w}.jpg"
                            convert_heic_to_jpg(photo_path, temp_jpg)

                            l, t, w_val, h_val = COORD_MAP[slide_idx]
                            slide.shapes.add_picture(temp_jpg, l, t, w_val, h_val)
                            os.remove(temp_jpg)

                    save_path = os.path.join(output_folder, f"{name.replace(' ', '_')}.pptx")
                    prs.save(save_path)
                except Exception as e:
                    st.error(f"Error for {name}: {e}")
                progress_bar.progress((index + 1) / len(df))
            status_text.success(f"🎉 Done! {len(df)} presentations created.")
