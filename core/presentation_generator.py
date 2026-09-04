import os
import pandas as pd
import streamlit as st
from pptx import Presentation
from concurrent.futures import ProcessPoolExecutor, as_completed
import multiprocessing

from core.config import COORD_MAP
from core.utils.portfolio_helpers import convert_heic_to_jpg, replace_text_recursive
from core.utils.google_drive import (
    get_all_photos_from_weeks,
    get_all_photos_from_weeks_oauth,
    download_public_file,
    get_drive_service,
    create_drive_folder,
    upload_drive_file
)

def remove_photo_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.type == 18:
            element = shape._element
            element.getparent().remove(element)

def process_single_student(index, row, mapping, root_path, folder_id, api_key, save_destination, global_theme, template_path, files_map, files_map_lower, stop_event, num_weeks, client_config):
    """
    Worker function to generate a single student's presentation.
    Returns a dictionary with the result and logs.
    """
    logs = []
    name = str(row[mapping["name"]]).strip()

    try:
        if stop_event.is_set():
            return {"name": name, "success": False, "logs": ["🛑 Stopped by user."], "error": "Stopped"}

        # Setup Drive Service if needed
        drive_service = None
        if folder_id:
            # If we have a folder_id and no API key, we assume OAuth or Service Account is needed
            if not api_key:
                try:
                    drive_service = get_drive_service(client_config)
                except Exception as e:
                    return {"name": name, "success": False, "logs": [f"❌ Drive Auth failed: {e}"], "error": str(e)}

        # Template and text replacement
        presentation = Presentation(template_path)
        student_class = str(row[mapping["class"]])
        section = str(row[mapping["section"]])
        theme_column = mapping["theme"]
        theme = str(row[theme_column]) if theme_column and pd.notna(row[theme_column]) else global_theme

        replacements = {"{{NAME}}": name, "{{CLASS}}": student_class, "{{SECTION}}": section, "{{THEME}}": theme}
        replace_text_recursive(presentation.slides[1].shapes, replacements)

        for week in range(1, num_weeks + 1):
            if stop_event.is_set():
                return {"name": name, "success": False, "logs": ["🛑 Stopped by user."], "error": "Stopped"}

            slide_index = week + 2
            slide = presentation.slides[slide_index - 1]
            remove_photo_placeholders(slide)
            expected_filename = f"{name}_W{week}.heic"

            photo_path = None
            if drive_service:
                file_id = files_map.get(expected_filename) or files_map_lower.get(expected_filename.lower())
                if file_id:
                    temp_heic = f"temp_{index}_{week}.heic"
                    logs.append(f"⏳ {name} W{week}: Downloading (OAuth)...")
                    import io
                    from googleapiclient.http import MediaIoBaseDownload
                    try:
                        request = drive_service.files().get_media(fileId=file_id)
                        with io.FileIO(temp_heic, 'wb') as fh:
                            downloader = MediaIoBaseDownload(fh, request)
                            done = False
                            while not done:
                                _, done = downloader.next_chunk()
                        photo_path = temp_heic
                    except Exception as e:
                        logs.append(f"❌ {name} W{week}: OAuth Download failed ({e})")

            elif folder_id and api_key:
                file_id = files_map.get(expected_filename) or files_map_lower.get(expected_filename.lower())
                if file_id:
                    temp_heic = f"temp_{index}_{week}.heic"
                    logs.append(f"⏳ {name} W{week}: Downloading (API Key)...")
                    success, error = download_public_file(file_id, api_key, temp_heic)
                    if success:
                        photo_path = temp_heic
                    else:
                        logs.append(f"❌ {name} W{week}: Download failed ({error})")

            elif root_path:
                local_path = os.path.join(root_path, f"Week {week}", expected_filename)
                if os.path.exists(local_path):
                    photo_path = local_path
                else:
                    logs.append(f"❌ {name} W{week}: Not found locally.")

            if photo_path:
                temp_jpg = f"temp_{index}_{week}.jpg"
                logs.append(f"⚙️ {name} W{week}: Converting HEIC to JPG...")
                if convert_heic_to_jpg(photo_path, temp_jpg):
                    left, top, width, height = COORD_MAP[slide_index]
                    slide.shapes.add_picture(temp_jpg, left, top, width, height)
                    logs.append(f"✅ {name} W{week}: Added to slide.")
                else:
                    logs.append(f"❌ {name} W{week}: Conversion failed.")

                if os.path.exists(temp_jpg):
                    os.remove(temp_jpg)
                if photo_path.startswith("temp_") and os.path.exists(photo_path):
                    os.remove(photo_path)

        # SAVE AND UPLOAD
        filename = f"{name.replace(' ', '_')}.pptx"
        temp_pptx = f"temp_{index}.pptx"
        presentation.save(temp_pptx)

        if save_destination == "Google Drive":
            if not drive_service:
                drive_service = get_drive_service(client_config)

            logs.append(f"📤 {name}: Uploading to Drive...")
            query = f"name = 'Finished_PPTs' and '{folder_id}' in parents and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
            results = drive_service.files().list(q=query, fields="files(id)", supportsAllDrives=True, includeItemsFromAllDrives=True).execute()
            files = results.get('files', [])

            if files:
                output_folder_id = files[0]['id']
            else:
                output_folder_id = create_drive_folder(drive_service, "Finished_PPTs", folder_id)

            if output_folder_id:
                success, err = upload_drive_file(drive_service, temp_pptx, output_folder_id, filename)
                if success:
                    logs.append(f"✅ {name}: Uploaded successfully.")
                else:
                    logs.append(f"❌ {name}: Upload failed ({err})")
            else:
                logs.append(f"❌ {name}: Could not find or create output folder.")

            if os.path.exists(temp_pptx):
                os.remove(temp_pptx)
        else:
            output_folder = os.path.join(root_path, "Finished_PPTs") if root_path else "Finished_PPTs"
            if not os.path.exists(output_folder):
                os.makedirs(output_folder, exist_ok=True)

            save_path = os.path.join(output_folder, filename)
            if os.path.exists(save_path):
                os.remove(save_path)
            os.rename(temp_pptx, save_path)
            logs.append(f"✅ {name}: Saved locally.")

        return {"name": name, "success": True, "logs": logs, "error": None}

    except Exception as error:
        return {"name": name, "success": False, "logs": logs + [f"💥 Critical Error - {error}"], "error": str(error)}


def generate_presentations(root_path, folder_id, api_key, drive_service, save_destination, template_path, dataframe, mapping, global_theme, num_weeks):
    # 1. Determine Output Destination (Local Only for initialization)
    if save_destination != "Google Drive":
        output_folder = os.path.join(root_path, "Finished_PPTs") if root_path else "Finished_PPTs"
        if not os.path.exists(output_folder):
            os.makedirs(output_folder, exist_ok=True)
    elif not drive_service:
        st.error("Drive service not authenticated. Please connect to Google Drive first.")
        return

    # Get client config for parallel workers (since st.secrets isn't available in subprocesses)
    client_config = {}
    try:
        client_config = st.secrets.get("google_auth", {})
    except Exception:
        pass

    progress_bar = st.progress(0)
    status_text = st.empty()
    log_area = st.expander("Detailed Generation Log", expanded=True)

    # Pre-fetch file list
    files_map = {}
    if drive_service:
        st.info("Scanning Google Drive (OAuth) sub-folders for photos...")
        files_map = get_all_photos_from_weeks_oauth(drive_service, folder_id)
    elif folder_id and api_key:
        st.info("Scanning Google Drive (API Key) sub-folders for photos...")
        files_map = get_all_photos_from_weeks(folder_id, api_key)

    files_map_lower = {k.lower(): v for k, v in files_map.items()} if files_map else {}

    # Parallel Execution Setup
    manager = multiprocessing.Manager()
    stop_event = manager.Event()

    # Collect tasks
    tasks = []
    for index, row in dataframe.iterrows():
        tasks.append((
            index, row, mapping, root_path, folder_id, api_key,
            save_destination, global_theme, template_path,
            files_map, files_map_lower, stop_event, num_weeks, client_config
        ))

    completed_count = 0
    total_students = len(dataframe)

    with ProcessPoolExecutor() as executor:
        # Map tasks to executor
        future_to_student = {executor.submit(process_single_student, *task): task[1][mapping["name"]] for task in tasks}

        try:
            for future in as_completed(future_to_student):
                student_name = future_to_student[future]
                result = future.result()

                completed_count += 1
                status_text.text(f"Processed {completed_count}/{total_students}: {result['name']}...")

                for log in result['logs']:
                    log_area.write(log)

                if not result['success']:
                    st.error(f"Error for {result['name']}: {result['error']}")

                progress_bar.progress(completed_count / total_students)

                # Check if user clicked STOP in the main thread
                if st.session_state.get("stop_generation", False):
                    stop_event.set()
        except Exception as e:
            st.error(f"Parallel execution error: {e}")

    if not st.session_state.get("stop_generation", False):
        status_text.success(f"🎉 Done! {total_students} presentations created.")
    else:
        status_text.info("Generation process was stopped.")


def render_generation(root_path, folder_id, api_key, drive_service, save_destination, uploaded_template, uploaded_data, dataframe, mapping, global_theme, num_weeks):
    st.subheader("Step 3: Generate Presentations")

    if "stop_generation" not in st.session_state:
        st.session_state.stop_generation = False

    col1, col2 = st.columns([3, 1])

    with col1:
        generate_btn = st.button("🔥 GENERATE ALL PRESENTATIONS", type="primary", use_container_width=True)

    with col2:
        if st.button("🛑 STOP", use_container_width=True):
            st.session_state.stop_generation = True
            st.rerun()

    if generate_btn:
        st.session_state.stop_generation = False
        if not uploaded_template or not uploaded_data or (not root_path and not folder_id) or dataframe is None:
            st.error("Missing requirements!")
            return
        if folder_id and not api_key and not drive_service:
            st.error("Please provide your Google API Key or authenticate via OAuth!")
            return

        template_path = "temp_template.pptx"
        with open(template_path, "wb") as template_file:
            template_file.write(uploaded_template.getbuffer())
        generate_presentations(root_path, folder_id, api_key, drive_service, save_destination, template_path, dataframe, mapping, global_theme, num_weeks)
